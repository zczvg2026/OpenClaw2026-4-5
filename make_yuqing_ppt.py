#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# Color palette
DARK_BLUE = RGBColor(0x0A, 0x2A, 0x5E)   # 深蓝
TEAL = RGBColor(0x00, 0x9E, 0xD0)          # 科技青
ORANGE = RGBColor(0xF5, 0x7C, 0x00)        # 橙色点缀
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MID_GRAY = RGBColor(0x7F, 0x8C, 0x96)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox

def add_title_bar(slide, title, subtitle=None):
    # Top blue bar
    add_rect(slide, 0, 0, 13.33, 1.4, DARK_BLUE)
    # Teal accent line
    add_rect(slide, 0, 1.4, 13.33, 0.06, TEAL)
    # Title text
    add_textbox(slide, 0.5, 0.2, 12, 0.9, title,
                font_size=36, font_color=WHITE, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, 0.5, 1.0, 12, 0.4, subtitle,
                    font_size=14, font_color=MID_GRAY, bold=False)

# ─────────────────────────────────────────────
# Slide 1: 封面
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
# Teal geometric accent - top right
add_rect(slide, 9.5, 0, 3.83, 0.15, TEAL)
add_rect(slide, 11.5, 0, 1.83, 0.15, ORANGE)
# Left accent bars
add_rect(slide, 0, 2.8, 0.12, 1.2, TEAL)
# Main title
add_textbox(slide, 0.8, 2.0, 10, 1.0, "域驍", font_size=72, bold=True,
            font_color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide, 0.8, 3.1, 10, 0.6, "AI原生全域增长引擎", font_size=36,
            font_color=TEAL, bold=True)
add_textbox(slide, 0.8, 3.85, 10, 0.5, "AI原生商业生态平台", font_size=20,
            font_color=MID_GRAY)
# Divider
add_rect(slide, 0.8, 4.5, 4.0, 0.04, TEAL)
# Bottom info
add_textbox(slide, 0.8, 4.7, 8, 0.4, "广州市太擎智能科技有限公司", font_size=14,
            font_color=MID_GRAY)
add_textbox(slide, 0.8, 5.1, 8, 0.4, "全域增长  ·  智能协同  ·  生态赋能", font_size=12,
            font_color=MID_GRAY)
# Brand slogan bottom right
add_textbox(slide, 8.5, 6.5, 4.5, 0.5, "系统有效  增长可靠", font_size=14,
            font_color=TEAL, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
# Slide 2: 核心定位
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "核心定位", "Domain Prime — AI-Native Business Ecosystem Platform")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

# Central tagline box
add_rect(slide, 0.5, 2.0, 12.33, 1.4, DARK_BLUE)
add_textbox(slide, 0.8, 2.15, 11.8, 1.1,
    "域驍是AI原生商业生态平台，生产AI产品和智能体，\n为企业和创业者双向赋能。",
    font_size=22, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)

# 4 keyword cards
keywords = [
    ("AI原生", "技术底座，不是传统AI应用", "科技感、时代感", TEAL),
    ("商业", "落地应用，面向真实商业场景", "实用性、确定性", ORANGE),
    ("生态", "完整生态，多向赋能", "循环效应、规模感", RGBColor(0x27, 0xAE, 0x60)),
    ("平台", "能力完整，主动创造", "规模化、可复制", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (word, desc, val, color) in enumerate(keywords):
    x = 0.5 + i * 3.15
    add_rect(slide, x, 3.7, 2.9, 3.3, WHITE)
    add_rect(slide, x, 3.7, 2.9, 0.12, color)
    add_textbox(slide, x+0.15, 3.9, 2.6, 0.6, word,
                font_size=28, font_color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, x+0.3, 4.55, 2.3, 0.03, color)
    add_textbox(slide, x+0.15, 4.7, 2.6, 1.0, desc,
                font_size=13, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.15, 6.0, 2.6, 0.5, val,
                font_size=12, font_color=color, bold=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# Slide 3: 八字定位
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "八字定位解析", "四维定位，全面诠释域驍")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

# Big 8-character display
add_rect(slide, 1.5, 1.9, 10.33, 1.2, DARK_BLUE)
add_textbox(slide, 1.5, 1.95, 10.33, 1.1,
    "四层架构 · 多向赋能 · 全域增长", font_size=32, bold=True,
    font_color=TEAL, align=PP_ALIGN.CENTER)

# 4 quadrants
quads = [
    ("AI原生", "技术底座，不是传统AI应用", TEAL),
    ("商业", "落地应用，面向真实商业场景", ORANGE),
    ("生态", "完整生态，多向赋能", RGBColor(0x27, 0xAE, 0x60)),
    ("平台", "能力完整，主动创造", RGBColor(0x8E, 0x44, 0xAD)),
]
positions = [(0.5, 3.5), (6.9, 3.5), (0.5, 5.6), (6.9, 5.6)]
for (word, desc, color), (x, y) in zip(quads, positions):
    add_rect(slide, x, y, 5.93, 1.8, WHITE)
    add_rect(slide, x, y, 0.1, 1.8, color)
    add_textbox(slide, x+0.25, y+0.2, 5.4, 0.5, word,
                font_size=24, font_color=color, bold=True)
    add_textbox(slide, x+0.25, y+0.75, 5.4, 0.9, desc,
                font_size=14, font_color=DARK_GRAY)

# ─────────────────────────────────────────────
# Slide 4: 四层架构
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "四层架构体系", "基座层 → 能力层 → 应用层 → 生态层")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

layers = [
    ("第四层", "多向生态层", "品牌企业 ↔ OPC创业者 ↔ AI应用 ↔ 解决方案", ORANGE, "循环赋能"),
    ("第三层", "应用层", "AI原生工具 + Agent集群解决方案\n招商型企业 | 创业者 | Agent集群", DARK_BLUE, "产品输出"),
    ("第二层", "能力层", "太驍技术平台 + OpenClaw + 物理硬件\n智能体开发 | 工作流编排 | 本地部署", RGBColor(0x27, 0xAE, 0x60), "技术支撑"),
    ("第一层", "基座层", "旷湖大数据平台 + 通用大模型\n1.8亿企业 | 毫秒响应 | 多模型融合", TEAL, "数据+AI底座"),
]

for i, (num, name, desc, color, tag) in enumerate(layers):
    y = 1.8 + i * 1.38
    add_rect(slide, 0.5, y, 12.33, 1.2, WHITE)
    add_rect(slide, 0.5, y, 0.1, 1.2, color)
    add_textbox(slide, 0.7, y+0.1, 1.5, 0.4, num, font_size=11, font_color=MID_GRAY)
    add_textbox(slide, 0.7, y+0.4, 2.0, 0.6, name, font_size=22, font_color=color, bold=True)
    add_textbox(slide, 2.8, y+0.2, 9.0, 0.8, desc, font_size=13, font_color=DARK_GRAY)
    add_textbox(slide, 11.5, y+0.3, 1.5, 0.6, tag, font_size=11, font_color=color,
                bold=True, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────
# Slide 5: 基座层详解
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "基座层详解", "旷湖大数据平台 + 通用大模型层")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

# Left: 旷湖
add_rect(slide, 0.5, 1.8, 5.8, 5.3, WHITE)
add_rect(slide, 0.5, 1.8, 5.8, 0.8, TEAL)
add_textbox(slide, 0.7, 1.9, 5.4, 0.6, "旷湖大数据平台", font_size=22,
            font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
items_left = [
    "中国：1.8亿企业 + 4亿关键决策人",
    "海外：3亿企业 + 13亿高管",
    "维度：1.5W+维度（国内）",
    "维度：1.5W+维度（海外）",
    "技术：毫秒级响应、PB级吞吐",
    "技术：实时图计算",
    "合规：合法可审计",
    "合规：标准制定参与者",
]
for i, item in enumerate(items_left):
    y = 2.75 + i * 0.48
    add_rect(slide, 0.7, y+0.05, 0.06, 0.3, TEAL)
    add_textbox(slide, 0.9, y, 5.2, 0.45, item, font_size=13, font_color=DARK_GRAY)

# Right: 大模型
add_rect(slide, 7.0, 1.8, 5.8, 5.3, WHITE)
add_rect(slide, 7.0, 1.8, 5.8, 0.8, ORANGE)
add_textbox(slide, 7.2, 1.9, 5.4, 0.6, "通用大模型层", font_size=22,
            font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
models = ["通义千问", "豆包", "智谱GLM", "Kimi", "Minimax", "……更多"]
for i, m in enumerate(models):
    x = 7.2 + (i % 3) * 1.85
    y = 2.85 + (i // 3) * 1.1
    add_rect(slide, x, y, 1.7, 0.85, DARK_BLUE)
    add_textbox(slide, x, y+0.1, 1.7, 0.65, m, font_size=14,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 7.2, 5.1, 5.4, 1.5,
    "多模型融合，适配不同场景需求\n提供AI智能能力底座",
    font_size=13, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# Slide 6: 生态循环
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "生态循环价值", "多位面协同共生，形成闭环")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

# Cycle flow
nodes = [
    ("品牌企业加入", DARK_BLUE, 0.5, 2.0),
    ("场景丰富度提升", TEAL, 3.5, 2.0),
    ("AI应用数量增加", ORANGE, 6.5, 2.0),
    ("OPC创业者选择更多", RGBColor(0x27,0xAE,0x60), 9.5, 2.0),
    ("创业成功案例增加", RGBColor(0x8E,0x44,0xAD), 6.5, 4.0),
    ("品牌企业增长动力增强", DARK_BLUE, 3.5, 4.0),
]
for label, color, x, y in nodes:
    add_rect(slide, x, y, 2.8, 1.1, color)
    add_textbox(slide, x, y+0.25, 2.8, 0.6, label, font_size=13,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrows as text labels (simplified)
arrow_texts = [
    ("→", 3.3, 2.2), ("→", 6.3, 2.2), ("→", 9.3, 2.2),
    ("↙", 6.3, 3.15), ("←", 9.3, 4.2), ("↖", 3.3, 3.15),
]
for sym, x, y in arrow_texts:
    add_textbox(slide, x, y, 0.5, 0.5, sym, font_size=20,
                font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Bottom value boxes
add_rect(slide, 0.5, 5.5, 3.8, 1.6, WHITE)
add_rect(slide, 0.5, 5.5, 3.8, 0.1, DARK_BLUE)
add_textbox(slide, 0.65, 5.7, 3.5, 0.4, "企业价值链", font_size=14,
            font_color=DARK_BLUE, bold=True)
add_textbox(slide, 0.65, 6.15, 3.5, 0.8,
    "尽早使用 → 沉淀数据 → 丰富生态\n→ 成本↓50% → 人效↑3倍 → 复购↑", font_size=11, font_color=DARK_GRAY)

add_rect(slide, 4.7, 5.5, 3.8, 1.6, WHITE)
add_rect(slide, 4.7, 5.5, 3.8, 0.1, TEAL)
add_textbox(slide, 4.85, 5.7, 3.5, 0.4, "创业者价值链", font_size=14,
            font_color=TEAL, bold=True)
add_textbox(slide, 4.85, 6.15, 3.5, 0.8,
    "学习AI → 选择项目 → 开始创业\n→ 自动运营 → 成功反哺 → 动力更强", font_size=11, font_color=DARK_GRAY)

add_rect(slide, 8.9, 5.5, 3.9, 1.6, WHITE)
add_rect(slide, 8.9, 5.5, 3.9, 0.1, ORANGE)
add_textbox(slide, 9.05, 5.7, 3.6, 0.4, "三位面协同", font_size=14,
            font_color=ORANGE, bold=True)
add_textbox(slide, 9.05, 6.15, 3.6, 0.8,
    "企业位面（需求方）\n创业者位面（参与方）\n平台位面（赋能方）", font_size=11, font_color=DARK_GRAY)

# ─────────────────────────────────────────────
# Slide 7: 核心竞争优势
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "核心竞争优势", "四大优势，构建竞争壁垒")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

advantages = [
    ("数据优势", TEAL, [
        "海量数据：1.8亿+4亿+3亿+13亿",
        "实时处理：毫秒级响应",
        "合规领先：标准制定参与者",
    ]),
    ("技术优势", DARK_BLUE, [
        "双平台架构：旷湖+太驍",
        "多模型融合：集成主流大模型",
        "低代码开发：快速迭代",
        "私有化部署：数据安全",
    ]),
    ("生态优势", ORANGE, [
        "多向循环：企业、创业者、平台相互赋能",
        "飞轮效应：越用越强",
        "规模效应：用户越多，价值越大",
    ]),
    ("产品优势", RGBColor(0x27,0xAE,0x60), [
        "AI原生：不是传统软件和工具",
        "场景聚焦：专注增长和提效",
        "标准化：可复制、可规模、可优化",
        "智能化：自动化、自增长、数据驱动",
    ]),
]
for i, (title, color, items) in enumerate(advantages):
    x = 0.5 + i * 3.15
    add_rect(slide, x, 1.8, 2.95, 5.3, WHITE)
    add_rect(slide, x, 1.8, 2.95, 0.7, color)
    add_textbox(slide, x, 1.9, 2.95, 0.5, title, font_size=20,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = 2.65 + j * 1.0
        add_rect(slide, x+0.15, y+0.08, 0.06, 0.3, color)
        add_textbox(slide, x+0.3, y, 2.5, 0.9, item, font_size=12, font_color=DARK_GRAY)

# ─────────────────────────────────────────────
# Slide 8: 价值主张
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.12, TEAL)
add_rect(slide, 0, 7.38, 13.33, 0.12, TEAL)
add_textbox(slide, 0.5, 0.4, 12.33, 0.8, "价值主张", font_size=36,
            font_color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_rect(slide, 0.5, 1.1, 1.5, 0.06, TEAL)

claims = [
    ("对企业", "用AI原生商业生态平台的工具，实现获客成本降低50%、人效提升3倍、复购率从20%→50%", TEAL),
    ("对创业者", "在AI原生商业生态平台，学会AI，用上工具，选对项目，创业成功。", ORANGE),
    ("对社会", "AI赋能，让每一个人都能拥有AI能力，实现个人和企业双向增长。", WHITE),
]
for i, (who, what, color) in enumerate(claims):
    y = 1.6 + i * 1.85
    add_rect(slide, 0.5, y, 12.33, 1.6, RGBColor(0x14, 0x35, 0x70))
    add_rect(slide, 0.5, y, 0.1, 1.6, color)
    add_textbox(slide, 0.8, y+0.15, 2.5, 0.5, who, font_size=20,
                font_color=color, bold=True)
    add_textbox(slide, 0.8, y+0.7, 11.8, 0.8, what, font_size=16, font_color=WHITE)

# ─────────────────────────────────────────────
# Slide 9: 成功指标
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "成功指标", "三维评估，量化生态价值")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

metrics = [
    ("企业端", DARK_BLUE, ["获客成本 ↓50%", "转化率 ↑显著", "复购率 20%→50%", "人效提升 3倍"]),
    ("创业者端", ORANGE, ["创业成功率 ↑", "收入增长显著", "项目存活率 ↑", "自动运营支撑"]),
    ("生态端", TEAL, ["用户数量 ↑", "应用数量 ↑", "案例数量 ↑", "付费率 ↑"]),
]
for i, (title, color, items) in enumerate(metrics):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 1.8, 3.95, 5.3, WHITE)
    add_rect(slide, x, 1.8, 3.95, 0.8, color)
    add_textbox(slide, x, 1.9, 3.95, 0.6, title, font_size=22,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = 2.8 + j * 1.0
        add_rect(slide, x+0.2, y+0.1, 0.06, 0.35, color)
        add_textbox(slide, x+0.4, y, 3.3, 0.7, item, font_size=16, font_color=DARK_GRAY, bold=True)

# ─────────────────────────────────────────────
# Slide 10: 愿景展望
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "愿景展望", "四大愿景，驱动未来增长")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

visions = [
    ("愿景1：连锁企业全域增长", DARK_BLUE, TEAL,
     "域驍招商版 + 智能运营中台（大龙虾）",
     "获客成本↓50% · 人效↑3倍 · 复购率20%→50%"),
    ("愿景2：个人创业者增长", ORANGE, RGBColor(0xFF,0xE0,0xB0),
     "域驍精英版 + 太驍智能体手机（小龙虾）",
     "月收入3千→1.9万 · 效率提升3倍"),
    ("愿景3：AI创业生态构建", RGBColor(0x27,0xAE,0x60), WHITE,
     "OPC联盟 + 培训体系 + 创业项目库",
     "创业成功率↑ · 创业周期↓"),
    ("愿景4：企业数字化转型", RGBColor(0x8E,0x44,0xAD), WHITE,
     "域驍企业版 + 私有化部署 + 定制开发",
     "数据统一 · 流程打通 · 智能决策"),
]
for i, (title, bg, fg, solution, effect) in enumerate(visions):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 1.8 + row * 2.7
    add_rect(slide, x, y, 6.1, 2.4, bg)
    add_textbox(slide, x+0.2, y+0.15, 5.7, 0.5, title, font_size=17,
                font_color=fg, bold=True)
    add_rect(slide, x+0.2, y+0.7, 5.7, 0.03, fg)
    add_textbox(slide, x+0.2, y+0.85, 5.7, 0.6, solution,
                font_size=13, font_color=WHITE)
    add_textbox(slide, x+0.2, y+1.5, 5.7, 0.6, effect,
                font_size=13, font_color=fg, bold=True)

# ─────────────────────────────────────────────
# Slide 11: 发展路径
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_title_bar(slide, "发展路径", "四阶段战略，稳步推进生态成熟")
add_rect(slide, 0, 1.46, 13.33, 0.06, TEAL)

stages = [
    ("第一阶段", "底座搭建", "✅ 已完成", DARK_BLUE, "旷湖+太驍平台建设"),
    ("第二阶段", "应用生产", "🔄 进行中", ORANGE, "AI工具批量生产"),
    ("第三阶段", "生态扩张", "📋 规划中", RGBColor(0x27,0xAE,0x60), "品牌企业引入，OPC联盟创建"),
    ("第四阶段", "生态成熟", "🔭 远期", RGBColor(0x8E,0x44,0xAD), "生态自我演化"),
]
# Timeline bar
add_rect(slide, 0.5, 4.0, 12.33, 0.08, MID_GRAY)
for i, (phase, name, status, color, desc) in enumerate(stages):
    x = 0.8 + i * 3.0
    add_rect(slide, 0.5, 4.0, 0.08, 0.08, color)
    # Circle node
    shape = slide.shapes.add_shape(9, Inches(x+0.3), Inches(3.7), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Top card
    add_rect(slide, x, 1.9, 2.7, 1.6, color)
    add_textbox(slide, x, 2.0, 2.7, 0.45, phase, font_size=13,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.45, 2.7, 0.45, name, font_size=18,
                font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.9, 2.7, 0.45, status, font_size=12,
                font_color=WHITE, align=PP_ALIGN.CENTER)
    # Bottom description
    add_textbox(slide, x, 4.35, 2.7, 0.6, desc, font_size=12,
                font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# Slide 12: 公司背景
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.12, TEAL)
add_rect(slide, 0, 7.38, 13.33, 0.12, TEAL)
add_textbox(slide, 0.5, 0.4, 12.33, 0.8, "公司背景", font_size=36,
            font_color=WHITE, bold=True)
add_rect(slide, 0.5, 1.1, 2.0, 0.06, TEAL)

# Left column
col1 = [
    ("母公司", "探迹科技", "全球AI独角兽"),
    ("市值", "13亿美元 ≈ 100亿人民币", "阿里巴巴、红杉、软银投资"),
    ("技术实力", "800+人技术团队", "服务6万家企业"),
    ("技术特点", "比钉钉早4个月发布仿生智能体", "企业级智能体解决方案提供商"),
]
for i, (label, val, sub) in enumerate(col1):
    y = 1.5 + i * 1.4
    add_rect(slide, 0.5, y, 5.8, 1.2, RGBColor(0x14,0x35,0x70))
    add_rect(slide, 0.5, y, 0.08, 1.2, TEAL)
    add_textbox(slide, 0.7, y+0.12, 5.4, 0.4, label, font_size=12, font_color=TEAL)
    add_textbox(slide, 0.7, y+0.48, 5.4, 0.4, val, font_size=15, font_color=WHITE, bold=True)
    add_textbox(slide, 0.7, y+0.85, 5.4, 0.35, sub, font_size=11, font_color=MID_GRAY)

# Right column
add_rect(slide, 6.8, 1.5, 6.0, 5.5, RGBColor(0x14,0x35,0x70))
add_textbox(slide, 7.0, 1.7, 5.6, 0.5, "品牌体系", font_size=18,
            font_color=TEAL, bold=True)
items_right = [
    "太驍智能科技 — 公司主体",
    "域驍 — 产品品牌（系列产品）",
    "真爱美家（003041）— 旗下上市公司",
    "（市值百亿人民币）",
    "",
    "旗下上市公司",
    "探迹科技 → 广州太驍智能科技",
    "（生态企业）",
]
for i, item in enumerate(items_right):
    if item:
        color = WHITE if i in [0,1,5] else MID_GRAY
        bold = i in [0,1,5]
        add_textbox(slide, 7.0, 2.3+i*0.5, 5.6, 0.45, item,
                    font_size=13 if bold else 12, font_color=color, bold=bold)

# ─────────────────────────────────────────────
# Slide 13: 结尾页
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
add_rect(slide, 0, 0, 0.15, 7.5, TEAL)
add_rect(slide, 13.18, 0, 0.15, 7.5, TEAL)
add_textbox(slide, 0, 2.2, 13.33, 1.0, "域驍", font_size=80, bold=True,
            font_color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.3, 13.33, 0.6, "重构商业  ·  赋能未来", font_size=28,
            font_color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_rect(slide, 4.5, 4.0, 4.33, 0.06, TEAL)
add_textbox(slide, 0, 4.3, 13.33, 0.5, "AI原生全域增长引擎", font_size=16,
            font_color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.9, 13.33, 0.4, "系统有效  增长可靠", font_size=14,
            font_color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 5.5, 13.33, 0.4, "广州市太驍智能科技有限公司", font_size=13,
            font_color=MID_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, 6.5, 13.33, 0.4, "全域增长  ·  智能协同  ·  生态赋能", font_size=12,
            font_color=MID_GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# 保存
# ─────────────────────────────────────────────
output_path = "/Users/mac/Desktop/zczvg/域驍介绍PPT_v1.0.pptx"
prs.save(output_path)
print(f"✅ PPT已保存: {output_path}")
