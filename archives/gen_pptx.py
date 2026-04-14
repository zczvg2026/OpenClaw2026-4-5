# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x0A, 0x0A, 0x1A)
ACCENT = RGBColor(0x00, 0x9E, 0xC8)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE8, 0xF0)
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x3A)
MID_BLUE = RGBColor(0x1A, 0x3A, 0x6A)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, left, top, width, height, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None:
        color = WHITE
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def add_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    tb(s, "AI Agent 时代来临", 0.8, 2.5, 11.5, 1.5, 54, True, WHITE, PP_ALIGN.CENTER)
    tb(s, "智能体技术正在重新定义人机协作的边界，开启自主决策的新纪元", 0.8, 4.2, 11.5, 0.8, 24, False, ACCENT, PP_ALIGN.CENTER)
    tb(s, "2026  |  Future of Work  |  域擎研究院", 0.8, 5.8, 11.5, 0.5, 16, False, LIGHT_GRAY, PP_ALIGN.CENTER)

def add_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    tb(s, "AGENDA", 1.0, 1.2, 3, 0.6, 18, True, ACCENT)
    tb(s, "目录", 1.0, 1.8, 4, 1.0, 44, True, WHITE)
    toc = [
        ("01 AI现状", "突破性进展与行业变革"),
        ("02 组织变革", "人机协同新模式"),
        ("03 域擎实践", "AI原生商业生态"),
        ("04 个人行动", "成为AI时代的OPC"),
    ]
    y = 3.4
    for num, desc in toc:
        tb(s, num, 1.0, y, 2.5, 0.6, 22, True, ORANGE)
        tb(s, desc, 3.8, y, 9, 0.6, 20, False, LIGHT_GRAY)
        y += 0.85

def add_slide(prs, title, lines, footer=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK_BG)
    bar = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()
    tb(s, title, 0.5, 0.18, 12, 0.7, 28, True, WHITE)
    y = 1.2
    for line in lines:
        size = 16 if len(line) > 60 else 18
        tb(s, line, 0.5, y, 12.3, 0.55, size, False, LIGHT_GRAY)
        y += 0.58
    if footer:
        tb(s, footer, 0.5, 6.85, 12, 0.4, 11, False, LIGHT_GRAY)

# --- Slides ---
add_cover(prs)
add_toc(prs)

add_slide(prs, "AI Agent 的突破性进展",
    [
        "AI Agent 智能体技术正经历从单点工具到自主系统的质的飞跃",
        "",
        "关键演进节点：",
        "  2022-2023：LLM突破，推理能力涌现，多步推理成为可能",
        "  2024：Agent Framework成熟，LangChain/AutoGPT引领潮流",
        "  2025：原生Agent架构，Claude Code/OpenClaw等产品涌现",
        "  2026：多Agent协作+垂直领域专业化，产业落地加速",
        "",
        "核心理念转变：从 AI辅助人 到 AI替代人执行",
    ], "域擎  |  AI Agent 时代来临  |  2026")

# Performance comparison slide
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
tb(s, "性能数据跨越式对比", 0.5, 0.25, 12, 0.65, 28, True, WHITE)
data = [
    ("GPT-4", "2023", "推理/编程基准", "83%"),
    ("Claude 3.5 Sonnet", "2024", "复杂任务处理", "91%"),
    ("Claude Code", "2025", "端到端编码Agent", "94%"),
    ("GPT-5 + Agent", "2026", "多Agent协作任务", "98%"),
]
y = 1.3
for model, year, task, score in data:
    pct = float(score.rstrip('%')) / 100.0
    bar = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(11.0 * pct), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tb(s, "{}  {}  |  {}  |  {}".format(model, year, task, score), 0.6, y + 0.1, 11.5, 0.55, 15, True, WHITE)
    y += 1.0
tb(s, "数据来源：第三方基准测试综合均值  |  域擎研究院整理  2026", 0.5, 6.85, 12, 0.4, 11, False, LIGHT_GRAY)

add_slide(prs, "商业应用落地开花",
    [
        "AI Agent已在多个行业实现规模化商业应用",
        "",
        "行业渗透率（2026Q1全球数据）：",
        "  软件开发 78%：Code Agent渗透率最高，头部企业全面采用",
        "  金融服务 54%：投研分析、风控、合规自动化",
        "  医疗健康 41%：病历整理、药物研发、患者随访",
        "  教育 38%：自适应学习、智能辅导、作业评估",
        "  制造业 29%：质检、工序优化、供应链调度",
        "",
        "注：65%企业已将AI Agent纳入核心业务流程（同比+31%）",
    ], "域擎  |  AI Agent 时代来临  |  2026")

# Stats slide
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
tb(s, "社会影响及未来展望", 0.5, 0.25, 12, 0.65, 28, True, WHITE)
stats = [
    ("-20%", "全球白领岗位\n预计减少比例"),
    ("2,643%", "AI相关岗位\n增长率（2024-2026）"),
    ("85%", "知识工作者\n使用AI工具比例"),
]
for i, (num, label) in enumerate(stats):
    x = 0.8 + i * 4.2
    box = s.shapes.add_shape(1, Inches(x), Inches(1.4), Inches(3.8), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = MID_BLUE
    box.line.color.rgb = ACCENT
    box.line.width = Pt(2)
    tb(s, num, x + 0.15, 1.7, 3.5, 1.1, 44, True, ORANGE, PP_ALIGN.CENTER)
    tb(s, label, x + 0.15, 2.9, 3.5, 0.9, 13, False, LIGHT_GRAY, PP_ALIGN.CENTER)
tb(s, "未来展望：AI Agent将重塑知识工作定义，重复性白领工作面临转型",
    0.5, 4.5, 12, 0.55, 16, False, WHITE)
tb(s, "新机遇：AI训练师、智能体架构师、人机协同设计师等新岗位大量涌现",
    0.5, 5.1, 12, 0.55, 16, False, ACCENT)
tb(s, "数据来源：麦肯锡、腾讯研究院、域擎研究院  |  2026",
    0.5, 6.85, 12, 0.4, 11, False, LIGHT_GRAY)

add_slide(prs, "组织变革：辅助协同进入AI时代",
    [
        "企业组织形态从金字塔向网络化+AI原生快速迭代",
        "",
        "传统组织  →  AI增强型  →  AI原生组织：",
        "  层级制减少：从5-7层压缩到3-4层",
        "  AI Agent替代中层管理的信息传递和协调职能",
        "  员工角色转变：决策者+AI监督者，而非执行者",
        "",
        "AI原生组织特征：",
        "  业务流程由Agent网络执行，人类负责规则和异常处理",
        "  绩效评估从完成多少任务，转向管理多少Agent协同产出",
    ], "域擎  |  AI Agent 时代来临  |  2026")

# Human-AI collab slide
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
tb(s, "HUMAN-AI COLLABORATION  人机协同新模式", 0.5, 0.2, 12, 0.65, 22, True, WHITE)
tb(s, "人类指挥一群AI员工，形成碳基决策+硅基执行的协作网络", 0.5, 0.9, 12, 0.5, 16, False, ACCENT)

roles = [
    ("人类角色：指挥官与编排师", ACCENT,
     "  战略规划与流程设计：专注顶层设计，制定业务目标与执行框架\n  协同监督与异常处理：负责关键决策审核与突发状况干预"),
    ("AI角色：硅基劳动力/数字员工", ORANGE,
     "  自主决策与多智能体协作\n  标准化业务执行：执行可复用、模块化的标准化任务"),
]
for i, (title, col, desc) in enumerate(roles):
    x = 0.5 + i * 6.3
    box = s.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(6.0), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = MID_BLUE
    box.line.color.rgb = col
    box.line.width = Pt(2)
    tb(s, title, x + 0.2, 1.65, 5.6, 0.5, 15, True, col)
    tb(s, desc, x + 0.2, 2.25, 5.6, 1.3, 13, False, LIGHT_GRAY)

tb(s, "协同增效实证数据", 0.5, 3.95, 12, 0.45, 15, True, WHITE)
sstats = [
    ("14-15%", "任务效率提升\n问题解决/小时"),
    ("26-50%", "产出提升\n代码提交/生产力"),
    ("55%", "处理能力提升\n客户支持吞吐量"),
    ("30-35%", "新手员工提升\nAI协助下成长"),
]
for i, (num, label) in enumerate(sstats):
    x = 0.5 + i * 3.1
    tb(s, num, x, 4.5, 2.9, 0.7, 28, True, ORANGE, PP_ALIGN.CENTER)
    tb(s, label, x, 5.25, 2.9, 0.8, 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)

tb(s, "新兴岗位：智能体编排师  |  首席智能体官  |  AI监督官  |  人机协同设计师",
    0.5, 6.7, 12, 0.5, 14, True, ACCENT)

# 域擎实践
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
tb(s, "域擎实践：AI原生商业生态", 0.5, 0.2, 12, 0.65, 28, True, WHITE)
tb(s, "从0到1构建智能体网络 · OPC联盟 × 垂直领域协议", 0.5, 0.9, 12, 0.5, 17, False, ACCENT)
practices = [
    ("OPC联盟平台", "垂直领域协议标准化",
     "制定行业智能体互操作协议\n连接企业Agent网络\n构建AI原生商业闭环"),
    ("财税智能体", "企业级AI数字员工",
     "自动化财税合规与申报\n7x24小时智能客服\n智能风控预警"),
    ("健康智能体", "S2B2C健康管理",
     "个人健康档案+AI分析\n医生-Agent协同诊断\n主动健康干预"),
]
for i, (name, sub, desc) in enumerate(practices):
    y = 1.55 + i * 1.8
    box = s.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(12.3), Inches(1.65))
    box.fill.solid()
    box.fill.fore_color.rgb = MID_BLUE
    box.line.fill.background()
    tb(s, name, 0.7, y + 0.12, 4, 0.5, 18, True, ORANGE)
    tb(s, sub, 0.7, y + 0.65, 4.5, 0.4, 13, False, ACCENT)
    tb(s, desc, 5.5, y + 0.12, 7, 1.4, 13, False, LIGHT_GRAY)
tb(s, "域擎研究院  |  AI Agent时代来临  |  2026", 0.5, 6.85, 12, 0.4, 11, False, LIGHT_GRAY)

# OPC行动
add_slide(prs, "个人行动：成为AI时代的OPC",
    [
        "OPC = Operator 操作者 + Player 玩家 + Creator 创造者",
        "",
        "三个身份：",
        "  Operator 操作者：会用AI工具，提升10倍个人生产力",
        "  Player 玩家：参与AI Agent网络，在协作中获取价值",
        "  Creator 创造者：用AI原生方式构建新产品、新流程、新商业模式",
        "",
        "第一步行动建议：",
        "  今天：用AI工具处理一件日常工作",
        "  本周：学习一个AI Agent协作框架",
        "  本月：找到自己的AI+专业交叉点",
    ], "域擎  |  AI Agent 时代来临  |  2026")

# 关键认知
add_slide(prs, "关键认知：为什么现在是关键节点",
    [
        "1. 技术成熟度拐点",
        "  2026年LLM推理成本下降90%，Agent调用成本接近零",
        "  多模态Agent已具备端到端任务完成能力",
        "",
        "2. 商业生态拐点",
        "  企业AI采购从试点转向全面部署",
        "  Agent网络效应显现：接入越多，价值越大",
        "",
        "3. 人才拐点",
        "  纯人类技能市场萎缩，AI协作能力成为核心竞争力",
        "  会用AI的人正在加速取代不会用AI的人",
    ], "域擎  |  AI Agent 时代来临  |  2026")

# 结语
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
tb(s, "AI不会取代人", 1.5, 1.8, 10, 1.2, 48, True, WHITE, PP_ALIGN.CENTER)
tb(s, "但会用AI的人会取代不会用AI的人", 1.5, 3.1, 10, 1.2, 36, True, ORANGE, PP_ALIGN.CENTER)
tb(s, "这不是选择，这是必然", 1.5, 4.4, 10, 0.7, 24, False, LIGHT_GRAY, PP_ALIGN.CENTER)
tb(s, "行动从今天开始", 1.5, 5.5, 10, 0.8, 22, True, ACCENT, PP_ALIGN.CENTER)
tb(s, "广州市太擎智能科技有限公司  |  域擎研究院", 1.5, 6.7, 10, 0.5, 14, False, LIGHT_GRAY, PP_ALIGN.CENTER)

# Save
out = "/Users/mac/Desktop/AI-Agent时代来临-域擎.pptx"
prs.save(out)
print("OK:" + out)
