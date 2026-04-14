#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE  = RGBColor(0x0A, 0x2A, 0x5E)
TEAL       = RGBColor(0x00, 0x9E, 0xD0)
ORANGE     = RGBColor(0xF5, 0x7C, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
DARK_GRAY  = RGBColor(0x2C, 0x3E, 0x50)
MID_GRAY   = RGBColor(0x7F, 0x8C, 0x96)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)
LIGHT_BLUE = RGBColor(0x14, 0x35, 0x70)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def R(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = "Microsoft YaHei UI"
    return tb

def bar(slide, title, subtitle=None):
    R(slide, 0, 0, 13.33, 1.4, DARK_BLUE)
    R(slide, 0, 1.4, 13.33, 0.06, TEAL)
    T(slide, 0.5, 0.18, 12, 0.85, title, size=34, color=WHITE, bold=True)
    if subtitle:
        T(slide, 0.5, 1.0, 12, 0.38, subtitle, size=13, color=MID_GRAY)

# ══════════════════════════════════════════════
# Slide A: 能力层详解
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
R(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
bar(s, "第二层：能力层", "太驍技术平台 + OpenClaw架构 + 物理硬件 — 构建技术能力和物理支撑")

# Section header
R(s, 0.5, 1.75, 12.33, 0.55, DARK_BLUE)
T(s, 0.7, 1.82, 12, 0.4, "核心价值：基于基座层，构建具体的技术能力和物理支撑", size=15, color=WHITE, bold=True)

cards = [
    ("太驍技术平台", TEAL, [
        "智能体开发",
        "知识库管理",
        "系统集成",
    ], "企业级安全"),
    ("OpenClaw", DARK_BLUE, [
        "工作流编排",
        "Agent 协同",
    ], "自动化能力"),
    ("物理硬件", ORANGE, [
        "大龙虾一体机",
        "小龙虾手机",
    ], "本地部署"),
]

for i, (name, color, items, value) in enumerate(cards):
    x = 0.5 + i * 4.2
    R(s, x, 2.55, 3.95, 4.55, WHITE)
    R(s, x, 2.55, 3.95, 0.08, color)
    R(s, x, 2.55, 0.12, 4.55, color)
    # Title
    T(s, x+0.2, 2.7, 3.55, 0.7, name, size=22, color=color, bold=True)
    # Divider
    R(s, x+0.2, 3.45, 3.55, 0.04, color)
    # Items
    for j, item in enumerate(items):
        y = 3.6 + j * 0.7
        R(s, x+0.3, y+0.12, 0.08, 0.3, color)
        T(s, x+0.5, y, 3.3, 0.6, item, size=15, color=DARK_GRAY)
    # Value badge
    R(s, x+0.2, 6.3, 3.55, 0.6, color)
    T(s, x+0.2, 6.38, 3.55, 0.45, f"→ {value}", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# Slide B: 应用层详解
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
R(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
bar(s, "第三层：应用层", "AI原生工具 + Agent集群解决方案 — 生产AI原生应用工具")

# Core value
R(s, 0.5, 1.75, 12.33, 0.55, DARK_BLUE)
T(s, 0.7, 1.82, 12, 0.4, "核心价值：基于能力层，生产AI原生应用工具和Agent集群解决方案", size=15, color=WHITE, bold=True)

# Three direction cards
directions = [
    ("招商型企业", ORANGE, [
        "全域招商增长解决方案",
        "为企业提供精准获客能力",
        "数据驱动增长引擎",
        "降低招商成本",
    ]),
    ("创业者", GREEN, [
        "增长解决方案输出",
        "创业项目赋能",
        "低成本启动路径",
        "自动化运营支撑",
    ]),
    ("Agent集群", PURPLE, [
        "多智能体协同",
        "复杂任务分解执行",
        "7×24持续运转",
        "可扩展智能体网络",
    ]),
]
for i, (name, color, items) in enumerate(directions):
    x = 0.5 + i * 4.2
    R(s, x, 2.55, 3.95, 4.55, WHITE)
    R(s, x, 2.55, 3.95, 0.08, color)
    R(s, x, 2.55, 0.12, 4.55, color)
    # Number badge
    R(s, x+0.15, 2.68, 0.5, 0.5, color)
    T(s, x+0.15, 2.72, 0.5, 0.4, f"{i+1}", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title
    T(s, x+0.75, 2.72, 3.0, 0.45, name, size=22, color=color, bold=True)
    # Divider
    R(s, x+0.2, 3.3, 3.55, 0.04, color)
    # Sub title
    T(s, x+0.2, 3.42, 3.55, 0.4, "三大应用方向之一", size=11, color=MID_GRAY)
    # Items
    for j, item in enumerate(items):
        y = 3.9 + j * 0.75
        R(s, x+0.3, y+0.1, 0.08, 0.3, color)
        T(s, x+0.5, y, 3.3, 0.65, item, size=14, color=DARK_GRAY)

# ══════════════════════════════════════════════
# Slide C: 生态层详解 — 7步循环
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
R(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
bar(s, "第四层：多向生态层", "品牌企业 ↔ OPC创业者 ↔ AI应用 ↔ 解决方案 — 形成闭环")

R(s, 0.5, 1.75, 12.33, 0.5, DARK_BLUE)
T(s, 0.7, 1.82, 12, 0.36, "核心价值：多位面协同共生，形成闭环，AI应用越来越丰富，生态越来越强大",
  size=14, color=WHITE, bold=True)

# 7-step cycle — horizontal flow
steps = [
    ("① 品牌企业加入",      DARK_BLUE),
    ("② 场景丰富度提升",    TEAL),
    ("③ AI应用数量增加",    ORANGE),
    ("④ OPC创业者选择更多", GREEN),
    ("⑤ 创业成功案例增加", PURPLE),
    ("⑥ 品牌企业增长增强", DARK_BLUE),
    ("⑦ 形成闭环",          TEAL),
]
for i, (label, color) in enumerate(steps):
    x = 0.3 + i * 1.83
    y = 2.55
    R(s, x, y, 1.68, 0.75, color)
    T(s, x, y+0.1, 1.68, 0.55, label, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 6:
        T(s, x+1.68, y+0.18, 0.15, 0.4, "→", size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Cycle return arrow
T(s, 0.3, 3.45, 12.73, 0.4, "↺  循环回到起点，持续迭代，生态不断壮大", size=15, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
R(s, 0.3, 3.85, 12.73, 0.06, TEAL)

# Three parties below
parties = [
    ("品牌企业", "需求方", "加入生态，提供场景和需求", DARK_BLUE),
    ("OPC创业者", "参与方", "选择项目，成功创业反哺生态", ORANGE),
    ("AI应用/解决方案", "双向引擎", "连接供需两侧，持续丰富", GREEN),
]
for i, (name, role, desc, color) in enumerate(parties):
    x = 0.5 + i * 4.2
    R(s, x, 4.15, 3.95, 3.0, WHITE)
    R(s, x, 4.15, 3.95, 0.08, color)
    R(s, x, 4.15, 0.12, 3.0, color)
    T(s, x+0.25, 4.25, 3.5, 0.5, name, size=18, color=color, bold=True)
    R(s, x+0.25, 4.8, 1.3, 0.45, color)
    T(s, x+0.25, 4.84, 1.3, 0.35, role, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(s, x+0.25, 5.35, 3.5, 1.5, desc, size=13, color=DARK_GRAY)

# ══════════════════════════════════════════════
# Slide D: 生态循环价值 — 三大价值链
# ══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
R(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
bar(s, "生态循环价值", "三位面协同，共生共长")

# 价值链1：企业
R(s, 0.5, 1.75, 12.33, 0.55, DARK_BLUE)
T(s, 0.7, 1.82, 12, 0.4, "企业价值链", size=16, color=TEAL, bold=True)

chain1 = [
    ("尽早使用", DARK_BLUE),
    ("沉淀数据", TEAL),
    ("丰富生态", ORANGE),
    ("成本↓50%", GREEN),
    ("人效↑3倍", PURPLE),
    ("复购20%→50%", DARK_BLUE),
]
for i, (node, color) in enumerate(chain1):
    x = 0.5 + i * 2.1
    R(s, x, 2.5, 1.9, 0.75, color)
    T(s, x, 2.6, 1.9, 0.55, node, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        T(s, x+1.9, 2.65, 0.2, 0.45, "→", size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 价值链2：创业者
R(s, 0.5, 3.55, 12.33, 0.55, ORANGE)
T(s, 0.7, 3.62, 12, 0.4, "创业者价值链", size=16, color=WHITE, bold=True)

chain2 = [
    ("学习AI", ORANGE),
    ("选择项目", DARK_BLUE),
    ("开始创业", TEAL),
    ("自动运营", GREEN),
    ("成功反哺", PURPLE),
    ("动力更强", ORANGE),
]
for i, (node, color) in enumerate(chain2):
    x = 0.5 + i * 2.1
    R(s, x, 4.3, 1.9, 0.75, color)
    T(s, x, 4.4, 1.9, 0.55, node, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        T(s, x+1.9, 4.45, 0.2, 0.45, "→", size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# 价值链3：三位面协同
R(s, 0.5, 5.3, 12.33, 0.5, TEAL)
T(s, 0.7, 5.37, 12, 0.36, "三位面协同", size=15, color=WHITE, bold=True)

triplets = [
    ("企业位面", "提供全域增长解决方案", "需求方", DARK_BLUE),
    ("创业者位面", "提供创业项目和增长方案", "参与方", ORANGE),
    ("平台位面", "提供AI工具和智能体集群", "赋能方", GREEN),
]
for i, (name, value, role, color) in enumerate(triplets):
    x = 0.5 + i * 4.2
    R(s, x, 6.0, 3.95, 1.2, WHITE)
    R(s, x, 6.0, 3.95, 0.08, color)
    R(s, x, 6.0, 0.12, 1.2, color)
    T(s, x+0.25, 6.08, 3.5, 0.4, name, size=15, color=color, bold=True)
    T(s, x+0.25, 6.5, 1.5, 0.35, value, size=11, color=DARK_GRAY)
    R(s, x+2.7, 6.08, 1.0, 0.35, color)
    T(s, x+2.7, 6.12, 1.0, 0.28, role, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════
output_path = "/Users/mac/Desktop/zczvg/域驍介绍PPT_补充页.pptx"
prs.save(output_path)
print(f"✅ 补充页已保存: {output_path}")
